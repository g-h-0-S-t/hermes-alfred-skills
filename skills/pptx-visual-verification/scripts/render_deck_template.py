"""Minimal runnable template for a render-to-PNG, vision-verified PPTX deck.

Usage: python render_deck_template.py
Produces assets/slides/slideNN.png (edit the slide functions) and deck.pptx.
Copy this file and extend the slide_* functions for a real deck.
"""
import os
from PIL import Image, ImageDraw, ImageFont

ASSETS = os.path.join(os.path.dirname(__file__), "..", "..", "assets")
FONT_DIR = os.path.join(ASSETS, "fonts")
SLIDE_DIR = os.path.join(ASSETS, "slides")
os.makedirs(SLIDE_DIR, exist_ok=True)
# Download Poppins if missing
for fn in ("Poppins-Bold.ttf", "Poppins-Regular.ttf"):
    p = os.path.join(FONT_DIR, fn)
    if not os.path.exists(p):
        os.makedirs(FONT_DIR, exist_ok=True)
        os.system(f"curl -sL -o {p} https://github.com/google/fonts/raw/main/ofl/poppins/{fn}")

W, H = 1600, 900
NAVY = (11, 46, 89); TEAL = (14, 155, 168); WHITE = (255, 255, 255)
PAPER = (247, 249, 252); CARD = (255, 255, 255); LINE = (223, 231, 240)
def F(p, s): return ImageFont.truetype(os.path.join(FONT_DIR, p), s)

def rr(d, box, r, fill, outline=None, ow=0):
    x0, y0, x1, y1 = box
    if x1 < x0: x1 = x0 + x1
    if y1 < y0: y1 = y0 + y1
    d.rounded_rectangle((x0, y0, x1, y1), r, fill=fill, outline=outline, width=ow)

def text(d, x, s, font, col, anchor="lm", line_h=None, max_w=None):
    if isinstance(x, (tuple, list)): x, y = x
    else: y = 0
    if max_w:
        wrapped = []
        for ln in s.split("\n"):
            cur = ""
            for wd in ln.split(" "):
                test = (cur + " " + wd).strip()
                if d.textlength(test, font=font) > max_w and cur: wrapped.append(cur); cur = wd
                else: cur = test
            wrapped.append(cur)
        s = "\n".join(wrapped)
    for i, ln in enumerate(s.split("\n")):
        bb = d.textbbox((0, 0), ln, font=font); tw = bb[2] - bb[0]; th = bb[3] - bb[1]
        cx, cy = x, y + i * (line_h or font.size + 6)
        if anchor in ("ma", "mm"): cx -= tw / 2
        if anchor in ("mm", "lm"): cy -= th / 2
        d.text((int(round(cx)), int(round(cy))), ln, font=font, fill=col)

def card(d, x, y, w, h, fill=CARD, line=None, shadow=True, r=18):
    if shadow: rr(d, (x + 7, y + 10, w, h), r, (205, 214, 226))
    rr(d, (x, y, w, h), r, fill, outline=line, ow=0 if line is None else 2)

def header(d, title, kicker=None):
    rr(d, (0, 0, W, 118), 0, NAVY); rr(d, (0, 118, W, 6), 0, TEAL)
    rr(d, (28, 34, 8, 52), 4, TEAL)
    text(d, (50, 38), title, F("Poppins-Bold.ttf", 40), WHITE, "lm")
    if kicker:
        kw = int(d.textlength(kicker, font=F("Poppins-Regular.ttf", 17))) + 28
        rr(d, (50, 82, kw, 28), 14, (28, 92, 140)); text(d, (64, 96), kicker, F("Poppins-Regular.ttf", 17), WHITE, "lm")

def footer(d, n):
    text(d, (44, H - 30), "Deck footer text", F("Poppins-Regular.ttf", 15), (110, 120, 130), "lm")
    text(d, (W - 60, H - 30), f"{n}/20", F("Poppins-Bold.ttf", 17), TEAL, "ma")

def s_title():
    img, d = __import__("PIL").Image.new("RGB", (W, H), NAVY), ImageDraw.Draw(__import__("PIL").Image.new("RGB", (W, H), NAVY))
    img = __import__("PIL").Image.new("RGB", (W, H), NAVY); d = ImageDraw.Draw(img)
    rr(d, (0, 560, W, 6), 0, TEAL)
    text(d, (44, 200), "Title Slide", F("Poppins-Bold.ttf", 62), WHITE, "lm")
    text(d, (44, 320), "Subtitle goes here", F("Poppins-SemiBold.ttf", 32), (191, 215, 245), "lm")
    return img

def s_content(n):
    img, d = __import__("PIL").Image.new("RGB", (W, H), PAPER), ImageDraw.Draw(__import__("PIL").Image.new("RGB", (W, H), PAPER))
    img = __import__("PIL").Image.new("RGB", (W, H), PAPER); d = ImageDraw.Draw(img)
    header(d, f"Content Slide {n}", "kicker")
    card(d, 44, 170, 720, 200, CARD, LINE)
    text(d, (72, 200), "Heading", F("Poppins-Bold.ttf", 26), TEAL, "lm")
    text(d, (72, 250), "Body text that wraps within the card.", F("Poppins-Regular.ttf", 21), (31, 41, 51), "lm", max_w=660)
    footer(d, n)
    return img

slides = [s_title()] + [s_content(i) for i in range(2, 21)]
for i, img in enumerate(slides, 1):
    img.save(os.path.join(SLIDE_DIR, f"slide{i:02d}.png"))
print("rendered", len(slides), "slides to", SLIDE_DIR)

# Assemble PPTX (option A: image slides)
from pptx import Presentation
from pptx.util import Inches
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
for i in range(1, len(slides) + 1):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(os.path.join(SLIDE_DIR, f"slide{i:02d}.png"), 0, 0, prs.slide_width, prs.slide_height)
out = os.path.join(os.path.dirname(__file__), "..", "..", "deck_example.pptx")
prs.save(out)
print("saved", out)
